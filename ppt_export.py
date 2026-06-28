import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
# Put downloaded PE_3F_KPI_Data.csv in same folder. This starter creates a title PPT.
prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[5]);
slide.shapes.title.text='PE-3F MSQCD Summary'
tx=slide.shapes.add_textbox(Inches(0.5),Inches(1.2),Inches(9),Inches(1)).text_frame
tx.text='Use dashboard CSV export as source. Extend this script for exact management slide layout.'
prs.save('PE_3F_KPI_Summary.pptx')
print('Created PE_3F_KPI_Summary.pptx')